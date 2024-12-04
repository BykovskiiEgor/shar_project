from pdfcreator.application.interfaces.presentation_reader import PresentationReader

from pptx import Presentation

class PresentationReader(PresentationReader):
    def read(self, file_path: str):
        prs = Presentation(file_path)
        
        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return text_runs