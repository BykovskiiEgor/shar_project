# from docx import Document

# def read_tables(file_path: str) -> dict:
#     doc = Document(file_path)
#     tables = doc.tables
    
#     all_tables_data = [] 
    
#     for table in tables:
#         table_data = []  
        
#         for row in table.rows:
#             row_data = [cell.text.strip() for cell in row.cells]
#             table_data.append(row_data)
        
#         all_tables_data.append(table_data)
    
#     return all_tables_data

# def save_tables_to_txt(file_path: str, tables_data: list, output_file: str) -> None:
#     with open(output_file, "w", encoding="utf-8") as file:
#         for table_index, table in enumerate(tables_data):
#             file.write(f"Table {table_index + 1}:\n")  
#             for row in table:
#                 file.write("\t".join(row) + "\n")  
#             file.write("\n")  

    
# output_file = "tables_output.txt"
# tables_data = read_tables('src/pdfcreator/test.docx')

# save_tables_to_txt('src/pdfcreator/test.docx', tables_data, output_file)    

from pptx import Presentation

def read(file_path: str):
    prs = Presentation(file_path)
    
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs
print(read('2024_10_Луховицкий_мукомольный_завод.pptx'))